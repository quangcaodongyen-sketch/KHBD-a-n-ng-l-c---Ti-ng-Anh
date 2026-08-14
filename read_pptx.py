import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text(shape):
    text = ""
    if hasattr(shape, "text"):
        text += shape.text + "\n"
    if shape.has_text_frame:
        text += shape.text_frame.text + "\n"
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text += cell.text_frame.text + " | "
            text += "\n"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            text += extract_text(child_shape)
    return text

def read_pptx(file_path):
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        print(f'--- Slide {i+1} ---')
        for shape in slide.shapes:
            try:
                print(extract_text(shape))
            except Exception as e:
                pass
        if slide.has_notes_slide:
            try:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    print(f"Notes: {notes}")
            except Exception:
                pass

if __name__ == '__main__':
    read_pptx('THCS.Dạy học TA lớp đa năng lực.pptx')
